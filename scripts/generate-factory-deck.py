#!/usr/bin/env python3
"""Generate Tejas-AI-Factory-Deck.pptx — run from repo root:
   scripts/.venv/bin/python scripts/generate-factory-deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(__file__).resolve().parent.parent / "docs" / "pitch" / "Tejas-AI-Factory-Deck.pptx"

NAVY = RGBColor(0x21, 0x29, 0x5C)
DEEP = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x02, 0x80, 0x90)
MINT = RGBColor(0x02, 0xC3, 0x9A)
CYAN = RGBColor(0x19, 0xC6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
MUTED = RGBColor(0x64, 0x74, 0x8B)
CHARCOAL = RGBColor(0x2D, 0x37, 0x48)
OFF = RGBColor(0xF4, 0xF7, 0xFA)
RED = RGBColor(0xFF, 0x33, 0x44)


def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_bar(slide, color=MINT):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.06))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def textbox(slide, left, top, width, height, text, size=18, bold=False, color=CHARCOAL, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def bullets(slide, items, left=0.6, top=1.8, width=8.8, size=15, color=CHARCOAL):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def dark_title(slide, title, sub="", foot=""):
    set_bg(slide, NAVY)
    add_bar(slide)
    textbox(slide, 0.6, 1.4, 8.8, 1.2, title, 36, True, WHITE)
    if sub:
        textbox(slide, 0.6, 2.5, 8.8, 0.8, sub, 18, False, ICE)
    if foot:
        textbox(slide, 0.6, 4.9, 8.8, 0.4, foot, 11, False, MINT)


def light_title(slide, title, sub=""):
    set_bg(slide, OFF)
    add_bar(slide)
    textbox(slide, 0.6, 0.45, 8.8, 0.7, title, 30, True, NAVY)
    if sub:
        textbox(slide, 0.6, 1.15, 8.8, 0.45, sub, 13, False, MUTED)


prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 1 TITLE
s = prs.slides.add_slide(prs.slide_layouts[6])
dark_title(s, "TEJAS AI", "The AI factory testing environment",
           "Train autonomous control · Stress the twin · Repair in VR · Deploy when ready")
textbox(s, 0.6, 3.2, 8.8, 0.5, "तेजस् · On-prem · Sovereign · Built for GPU halls", 14, False, CYAN)
notes(s, "Open: You cannot learn autonomous cooling on live H100s.")

# 2 PROBLEM
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "AI factories broke the old playbook", "Seconds to act · Weeks of jobs at risk")
bullets(s, [
    "Rack power spikes before inlet temp — you have seconds, not minutes",
    "One cooling trip = ₹2+ Cr lost compute + SLA penalties",
    "Phaidra proved autonomous control on NVIDIA-scale halls",
    "Everyone else still runs 1990s rule-based BMS logic",
], top=1.65)
box = s.shapes.add_shape(1, Inches(0.6), Inches(4.2), Inches(2.5), Inches(0.9))
box.fill.solid(); box.fill.fore_color.rgb = TEAL; box.line.fill.background()
textbox(s, 0.75, 4.3, 2.2, 0.35, "30–40%", 22, True, WHITE)
textbox(s, 0.75, 4.6, 2.2, 0.35, "of power = cooling", 10, False, ICE)
notes(s, "Mid-market has no way to practice before flipping the switch.")

# 3 CANNOT TEST LIVE
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "The real hall is not a lab")
bullets(s, [
    "❌  Live: failed job = customer churn          →  ✅  Twin: heat wave in one drag",
    "❌  Live: wrong setpoint = chip throttle       →  ✅  Twin: chiller trip — recover in sim",
    "❌  Live: no rollback for the board            →  ✅  Shadow: log would-do vs your BMS",
    "❌  Live: MoC / insurance nightmare            →  ✅  Trust ladder with instant revert",
], top=1.55, size=13)
textbox(s, 0.6, 4.35, 8.8, 0.6,
        '"Would you let an untested AI drive your GPUs on Thursday\'s 47°C day?"', 14, True, DEEP, PP_ALIGN.CENTER)
notes(s, "Fear frame — then offer sandbox.")

# 4 TESTING ENV
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "Tejas = your AI factory testing environment", "Full physics twin · Zero risk to silicon")
bullets(s, [
    "Calibrated digital twin — topology + historian, not a drawing toy",
    "Autonomous brain vs baseline — live PUE 1.3 vs 1.5 at 47°C",
    "Stress: heat waves · load spikes · chiller trips · GPU faults",
    "Graduate path: Twin → Shadow (read-only) → Advise → You take control",
], top=1.6)
notes(s, "We give a sandbox that looks like their hall.")

# 5 DEMO
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, NAVY)
add_bar(s)
textbox(s, 0.6, 1.8, 8.8, 0.8, "DEMO NOW", 40, True, MINT, PP_ALIGN.CENTER)
bullets(s, [
    "1. /datacenter — Autonomous Control · PUE live",
    "2. Weather → 47°C — baseline panics · Tejas holds",
    '3. "Is any machine going down?" → GPU-16 red',
    "4. Hand phone → /vr — stand inside the hall",
], top=2.7, size=14, color=ICE)
notes(s, "2 min demo. End on VR handoff.")

# 6 AUTONOMOUS
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "What the testing environment trains", "Feed-forward control — the AI factory requirement")
bullets(s, [
    "Acts on rack power + weather before inlet spikes (<10s)",
    "Explains every move — heat wave COP protection",
    "Dual run: Tejas vs dumb BMS — savings gap visible live",
    "Shadow on real hall: their historian proves ₹ before control",
], top=1.6)
notes(s, "Phaidra-class principle, mid-market access.")

# 7 TRUST LADDER
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "You graduate — not us", "From testing env to real hall on your schedule")
phases = ["Twin\n(test)", "Shadow\n(read-only)", "Advise\n(you apply)", "Supervised\n(confirm)", "Autonomous\n(you own)"]
for i, p in enumerate(phases):
    x = 0.35 + i * 1.88
    sh = s.shapes.add_shape(1, Inches(x), Inches(1.85), Inches(1.72), Inches(2.2))
    sh.fill.solid()
    sh.fill.fore_color.rgb = TEAL if i <= 2 else WHITE
    sh.line.color.rgb = TEAL
    textbox(s, x, 2.0, 1.72, 1.8, p, 11, True, WHITE if i <= 2 else NAVY, PP_ALIGN.CENTER)
textbox(s, 0.6, 4.5, 8.8, 0.4, "We hand you the keys when YOUR data says yes.", 13, True, DEEP)
notes(s, "Commercial: shadow low/free, savings-share after they choose.")

# 8 HERO VR
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DEEP)
add_bar(s, CYAN)
textbox(s, 0.6, 0.5, 8.8, 0.7, "Repair in VR", 34, True, WHITE)
textbox(s, 0.6, 1.15, 8.8, 0.45, "Less downtime · Less cost · Never walk blind", 16, False, ICE)
rows = [
    ("2:14 AM alarm", "GPU-16 flagged 14 months early"),
    ("Senior tech on a plane", "VR walk to exact rack — junior on site"),
    ("Open rack blind", "Exploded design — filter, fan tray, busbar"),
    ("₹2+ Cr outage", "25-min planned fix · parts pre-staged"),
]
for i, (old, new) in enumerate(rows):
    y = 1.75 + i * 0.85
    textbox(s, 0.6, y, 3.8, 0.35, old, 12, False, RGBColor(0xFF, 0xAA, 0xAA))
    textbox(s, 4.6, y, 5, 0.35, "→  " + new, 12, True, MINT)
notes(s, "EMOTIONAL PEAK. Put phone in their hand.")

# 9 VR FLOW
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "VR repair flow", "Walk in · Pick · Inspect · Fix · Blue")
bullets(s, [
    "Scan QR → stand inside live twin (/vr)",
    "Walk to GPU-16 — pulsing red · reticle pick",
    "Aspect lenses — cooling path: filter + fan tray",
    "Field AR on site — 4-step checklist · rack turns blue",
    "Desktop syncs — work order closed automatically",
], top=1.55, size=14)
notes(s, "Chennai tech practiced in VR before the outage.")

# 10 GPU-16
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "Predict in twin · Practice in VR · Fix on site")
bullets(s, [
    "Inlet 3.4°C hotter than load explains — 14-month creep",
    "Root cause: CRAC-2 filter + rear fan-tray bearing",
    "Work order + email with exact parts drafted automatically",
    "Caught in testing twin → validated on historian in shadow",
], top=1.6)
notes(s, 'Chat: "Is any machine going down?"')

# 11 NUMBERS
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "What the testing env proves", "12 MW GPU hall · before one setpoint hits silicon")
stats = [("₹1.8–3.2 Cr", "cooling saved/yr"), ("₹2+ Cr", "one avoided outage"),
         ("₹8.3L+", "per PM fix"), ("₹90–160 Cr", "50-hall fleet")]
for i, (v, l) in enumerate(stats):
    x = 0.5 + (i % 2) * 4.6
    y = 1.7 + (i // 2) * 1.35
    box = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.2), Inches(1.1))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ICE
    textbox(s, x + 0.15, y + 0.1, 3.9, 0.45, v, 22, True, NAVY)
    textbox(s, x + 0.15, y + 0.55, 3.9, 0.4, l, 11, False, MUTED)
notes(s, "Shadow puts THEIR historian on this slide.")

# 12 WHO
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "Built for halls Phaidra doesn't serve — yet")
bullets(s, [
    "GPU colo / edge AI DC — sovereign on-prem testing sandbox",
    "Enterprise AI lab — internal team owns control policy",
    "DC adding GPU block — test liquid + air mix before go-live",
    "ESCO / integrator — white-label testing env for customers",
], top=1.6)
notes(s, "Not for greenfield with zero telemetry.")

# 13 COMMERCIAL
s = prs.slides.add_slide(prs.slide_layouts[6])
light_title(s, "Pay for proof — not promises")
bullets(s, [
    "Twin build / testing env — paid on-ramp (credited to Y1)",
    "Shadow 30 days — read-only · would-do report · low / free",
    "Advisory + control — savings-share 20–40% when THEY graduate",
    "Land-and-expand: one hall → campus → 50-site fleet",
], top=1.6)
notes(s, "They flip the switch — we don't force autonomous.")

# 14 CLOSE
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, NAVY)
add_bar(s)
textbox(s, 0.6, 1.0, 8.8, 0.6, "Don't learn autonomous cooling on live GPUs.", 22, False, ICE)
textbox(s, 0.6, 1.65, 8.8, 1.4,
        "Learn it in the twin.\nRepair it in VR.\nOwn the real hall when you're ready.", 28, True, WHITE)
textbox(s, 0.6, 3.5, 8.8, 1.0,
        "30 days read-only on your historian.\nWe build your testing environment — and show you the rack you'd lose on a 47°C Thursday.",
        13, False, MINT)
notes(s, "CTA: calendar hold for shadow scoping call.")

prs.save(str(OUT))
print(f"Created {OUT}")